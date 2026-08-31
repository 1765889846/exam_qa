"""文档解析：PDF / Office / 纯文本 → ParsedDocument。"""

from __future__ import annotations

import base64
import json
import logging
import os
import re
import shutil
import subprocess
import tempfile
from dataclasses import dataclass, field
from pathlib import Path

from src.config import config
from src.exceptions import BadRequestException, UnsupportedFormatException

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pdf", ".txt", ".md", ".doc", ".docx", ".pptx"}


def _run_capture(args: list[str], *, timeout: float | None = None, encoding: str = "utf-8"):
    """subprocess.run 的 Windows 安全版：stdout/stderr 写临时文件而非管道。

    子进程若再派生子进程（COM dllhost / MinerU 子进程）会继承管道句柄，
    导致 communicate() 在超时后仍永久阻塞；临时文件无句柄继承问题。
    """
    import tempfile as _tf

    with _tf.TemporaryFile() as _out, _tf.TemporaryFile() as _err:
        proc = subprocess.run(args, stdout=_out, stderr=_err, timeout=timeout)
        _out.seek(0)
        _err.seek(0)
        stdout = _out.read().decode(encoding, errors="replace")
        stderr = _err.read().decode(encoding, errors="replace")
    return proc, stdout, stderr


@dataclass
class ParsedPage:
    page: int | None  # PDF/PPT 1-based；纯文本/docx 为 None
    text: str


@dataclass
class ParsedBlock:
    """MinerU 结构化 block：结构还原的最小单元。"""

    block_type: str  # text|title|table|formula|formula_inline|image|image_caption|table_caption|figure_caption|header|footer|...
    text: str  # 可检索文本（表格/公式/图片摘要已转文本）
    page: int | None  # 1-based
    level: int = 0  # 标题层级 1-6；非标题为 0
    html: str = ""  # 表格原始 HTML
    latex: str = ""  # 公式 LaTeX
    image_path: str = ""  # 图片原始文件路径（多模态摘要用）
    caption: str = ""  # 图片/表格说明
    order: int = 0  # 原始顺序（跨页排序用）
    bbox: tuple | None = None


@dataclass
class ParsedDocument:
    pages: list[ParsedPage]
    blocks: list[ParsedBlock] = field(default_factory=list)

    @property
    def full_text(self) -> str:
        return "\n\n".join(p.text for p in self.pages if p.text.strip())

    @property
    def has_blocks(self) -> bool:
        return bool(self.blocks)


def _page_num(meta: dict) -> int | None:
    raw = meta.get("page_number") or meta.get("page")
    if raw is None:
        return None
    try:
        return int(raw)
    except (TypeError, ValueError):
        return None


def _pages_from_pymupdf4llm(raw) -> ParsedDocument | None:
    if isinstance(raw, str):
        text = raw.strip()
        return ParsedDocument([ParsedPage(page=None, text=text)]) if text else None
    pages = []
    for chunk in raw:
        text = (chunk.get("text") or "").strip()
        if text:
            pages.append(ParsedPage(page=_page_num(chunk.get("metadata") or {}), text=text))
    return ParsedDocument(pages) if pages else None


def _pymupdf4llm_kwargs(*, force_ocr: bool | None = None) -> dict:
    p = config.parsing
    return {
        "page_chunks": True,
        "use_ocr": p.pdf_use_ocr,
        "force_ocr": p.pdf_force_ocr if force_ocr is None else force_ocr,
        "ocr_language": p.pdf_ocr_language,
    }


def _parse_pdf_pymupdf4llm(path: str, *, force_ocr: bool | None = None) -> ParsedDocument | None:
    import pymupdf4llm

    return _pages_from_pymupdf4llm(
        pymupdf4llm.to_markdown(path, **_pymupdf4llm_kwargs(force_ocr=force_ocr))
    )


def _parse_pdf_fitz(path: str) -> ParsedDocument:
    import fitz

    doc = fitz.open(path)
    try:
        pages = [
            ParsedPage(page=i, text=t)
            for i, page in enumerate(doc, 1)
            if (t := page.get_text().strip())
        ]
        return ParsedDocument(pages)
    finally:
        doc.close()


def _mineru_available(cmd: str) -> bool:
    """MinerU CLI 是否可用（PATH 可找到）。"""
    return bool(shutil.which(cmd))


def _pdf_kind(path: str) -> str:
    """抽样判断 PDF 类型：text（原生文本）| scanned（扫描件/图片型）| mixed。

    原生文本 PDF 直接提取文本；扫描件/混合型交给 OCR / MinerU。
    采样首页、次页、中间页、末页，避免只按首字符数误判。
    """
    try:
        import fitz

        doc = fitz.open(path)
        try:
            n = doc.page_count
            if n == 0:
                return "scanned"
            total = 0
            sampled = 0
            for i in sorted({0, 1, n // 2, n - 1}):
                if i >= n:
                    continue
                total += len(re.sub(r"\s", "", doc[i].get_text() or ""))
                sampled += 1
            if sampled == 0 or total < 30:
                return "scanned"
            avg = total / sampled
            if avg < 30:
                return "scanned"
            return "mixed" if avg < 200 else "text"
        finally:
            doc.close()
    except Exception:
        return "scanned"


def _html_table_to_text(html: str) -> str:
    """把 MinerU 的表格 HTML 转成可检索的纯文本（单元格用 | 分隔）。"""
    if not html:
        return ""
    text = re.sub(r"<t[dh][^>]*>", " | ", html, flags=re.I)
    text = re.sub(r"</t[r]>", "\n", text, flags=re.I)
    text = re.sub(r"</t[dh]>", "", text, flags=re.I)
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"&nbsp;?", " ", text, flags=re.I)
    lines = [re.sub(r"\s*[|]\s*$", "", ln).strip() for ln in text.splitlines()]
    return "\n".join(ln for ln in lines if ln.strip())


def _table_headers(html: str) -> str:
    """从表格 HTML 提取列头（<th>），写入结构化 metadata。"""
    if not html:
        return ""
    heads = re.findall(r"<t[h][^>]*>(.*?)</t[h]>", html, flags=re.I | re.S)
    cleaned = [" ".join(re.sub(r"<[^>]+>", " ", h).split()) for h in heads]
    return " | ".join(c for c in cleaned if c)


def _block_text(block: dict) -> str:
    """content_list 单个 block 的可检索文本（text/table/formula 按类型拼接）。"""
    btype = block.get("type") or ""
    if btype in ("formula", "formula_inline"):
        latex = block.get("latex") or ""
        return f"$${latex}$$" if latex else ""
    if btype == "table":
        text = block.get("text") or ""
        html = block.get("html") or ""
        return text if text else _html_table_to_text(html)
    text = block.get("text") or ""
    if text:
        return text
    parts = []
    for ln in block.get("lines") or []:
        if isinstance(ln, dict):
            t = ln.get("text") or ln.get("content") or ""
        else:
            t = str(ln)
        if t.strip():
            parts.append(t.strip())
    return "\n".join(parts)


_CAPTION_TYPES = frozenset({"image_caption", "table_caption", "figure_caption"})


def _caption_text(block: dict) -> str:
    """图片/表格说明文本（兼容 text 与 lines 两种形态）。"""
    text = (block.get("text") or "").strip()
    if text:
        return text
    parts = []
    for ln in block.get("lines") or []:
        if isinstance(ln, dict):
            t = ln.get("text") or ln.get("content") or ""
        else:
            t = str(ln)
        if t.strip():
            parts.append(t.strip())
    return " ".join(parts)


def _as_bbox(raw) -> tuple | None:
    if not isinstance(raw, (list, tuple)) or len(raw) != 4:
        return None
    try:
        return tuple(float(v) for v in raw)
    except (TypeError, ValueError):
        return None


def _resolve_image_path(img_path: str | None, base_dir: Path | None) -> str:
    """MinerU 图片路径可能是相对 json 输出目录的，解析为绝对路径。"""
    if not img_path:
        return ""
    p = Path(img_path)
    if not p.is_absolute() and base_dir is not None:
        p = base_dir / p
    return str(p.resolve()) if p.exists() else ""


def _blocks_to_pages(blocks: list[ParsedBlock]) -> list[ParsedPage]:
    """按页聚合 block 文本，保证 pages 视图与 blocks 一致。"""
    by_page: dict[int, list[str]] = {}
    for b in blocks:
        if not b.text:
            continue
        by_page.setdefault(b.page or 0, []).append(b.text)
    pages = []
    for k in sorted(by_page):
        pages.append(
            ParsedPage(page=None if k == 0 else k, text="\n\n".join(by_page[k]))
        )
    return pages


def _mineru_json_to_document(content: dict, *, base_dir: Path | None = None) -> ParsedDocument:
    """重组 MinerU content_list 为带 block 结构的 ParsedDocument（结构还原）。

    - 保留 title 层级、table HTML、formula、image 原图路径与说明；
    - 图片/表格说明（caption）就近挂载，无归属的保留为独立 caption block；
    - 页眉页脚保留在 blocks 中（结构还原），正文切片阶段再决定取舍。
    """
    items = content.get("content_list") or []
    raw: list[dict] = []
    for i, block in enumerate(items):
        if isinstance(block, dict):
            item = dict(block)
            item.setdefault("_seq", i)
            raw.append(item)

    # caption 挂载：优先挂到它前面的 image/table（MinerU 输出图在前、说明在后），
    # 无前驱时再找后面的，距离超过 3 个 block 视为无归属。
    targets = [
        (int(b.get("_seq", 0)), b.get("type"))
        for b in raw
        if b.get("type") in ("image", "table")
    ]
    caption_by_target: dict[int, str] = {}
    orphan_captions: list[tuple[int, int | None, str]] = []
    for block in raw:
        if block.get("type") not in _CAPTION_TYPES:
            continue
        text = _caption_text(block)
        if not text:
            continue
        seq = int(block.get("_seq", 0))
        candidates = [t for t in targets if t[0] < seq] or targets
        best_t, best_d = -1, 10**9
        for tseq, _ttype in candidates:
            d = abs(tseq - seq)
            if d < best_d:
                best_t, best_d = tseq, d
        if best_t >= 0 and best_d <= 3:
            caption_by_target[best_t] = text
        else:
            try:
                page = int(block.get("page_idx", 0)) + 1
            except (TypeError, ValueError):
                page = None
            orphan_captions.append((seq, page, text))

    blocks: list[ParsedBlock] = []
    for block in raw:
        try:
            page = int(block.get("page_idx", 0)) + 1
        except (TypeError, ValueError):
            page = None
        btype = str(block.get("type") or "text")
        seq = int(block.get("_seq", 0))
        order = block.get("order") if block.get("order") is not None else seq

        if btype in _CAPTION_TYPES:
            continue  # 已挂载或进入 orphan_captions
        if btype == "table":
            html = block.get("html") or ""
            text = (block.get("text") or "").strip() or _html_table_to_text(html)
            blocks.append(
                ParsedBlock(
                    block_type="table",
                    text=text,
                    page=page,
                    html=html,
                    caption=caption_by_target.get(seq, ""),
                    order=order,
                    bbox=_as_bbox(block.get("bbox")),
                )
            )
        elif btype == "image":
            blocks.append(
                ParsedBlock(
                    block_type="image",
                    text=(block.get("text") or "").strip(),
                    page=page,
                    image_path=_resolve_image_path(block.get("img_path"), base_dir),
                    caption=caption_by_target.get(seq, ""),
                    order=order,
                    bbox=_as_bbox(block.get("bbox")),
                )
            )
        elif btype in ("formula", "formula_inline"):
            latex = block.get("latex") or ""
            blocks.append(
                ParsedBlock(
                    block_type=btype,
                    text=f"$${latex}$$" if latex else "",
                    page=page,
                    latex=latex,
                    order=order,
                    bbox=_as_bbox(block.get("bbox")),
                )
            )
        elif btype == "title":
            try:
                level = int(block.get("level") or 1)
            except (TypeError, ValueError):
                level = 1
            blocks.append(
                ParsedBlock(
                    block_type="title",
                    text=(block.get("text") or "").strip(),
                    page=page,
                    level=level,
                    order=order,
                    bbox=_as_bbox(block.get("bbox")),
                )
            )
        else:
            blocks.append(
                ParsedBlock(
                    block_type=btype,
                    text=_block_text(block).strip(),
                    page=page,
                    order=order,
                    bbox=_as_bbox(block.get("bbox")),
                )
            )

    for seq, page, text in orphan_captions:
        blocks.append(ParsedBlock(block_type="caption", text=text, page=page, order=seq))

    blocks.sort(key=lambda b: (b.page or 0, b.order))
    return ParsedDocument(pages=_blocks_to_pages(blocks), blocks=blocks)


def _mineru_json_to_pages(content: dict) -> ParsedDocument:
    """兼容入口：仅按页重组文本（无结构信息）。"""
    return _mineru_json_to_document(content)


def _mineru_md_to_pages(md_text: str) -> ParsedDocument:
    """无 JSON 时的回退：Markdown 整体作为一页。"""
    text = md_text.strip()
    return ParsedDocument([ParsedPage(page=None, text=text)] if text else [])


def _image_mime(path: Path) -> str:
    return {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".webp": "image/webp",
        ".gif": "image/gif",
        ".bmp": "image/bmp",
    }.get(path.suffix.lower(), "")


def _summarize_image(image_path: str, caption: str, page: int | None) -> str | None:
    """用配置的视觉模型生成图片/图表摘要（多模态视觉理解）。

    未配置 VISUAL_MODEL / 图片缺失 / 调用失败时返回 None，由调用方回退占位，
    不阻断入库主链路。
    """
    p = config.parsing
    if not p.visual_model or not image_path:
        return None
    path = Path(image_path)
    if not path.exists():
        logger.warning("图片不存在，跳过视觉摘要: %s", image_path)
        return None
    mime = _image_mime(path)
    if not mime:
        logger.warning("不支持的图片格式，跳过视觉摘要: %s", path.name)
        return None
    base_url = p.visual_base_url or config.llm.base_url or "https://api.openai.com/v1"
    api_key = p.visual_api_key or config.llm.api_key
    if not api_key:
        logger.warning("VISUAL_MODEL 已配置但无 VISUAL_API_KEY / LLM_API_KEY，跳过视觉摘要")
        return None
    try:
        data_url = (
            f"data:{mime};base64,"
            f"{base64.b64encode(path.read_bytes()).decode('ascii')}"
        )
    except OSError as e:
        logger.warning("读取图片失败，跳过视觉摘要: %s", e)
        return None

    prompt = (
        "请用中文描述这张教学资料图片/图表的核心内容与关键信息，"
        "2-3 句话，适合作为检索摘要，不要输出多余内容。"
    )
    if caption:
        prompt += f"\n图片说明（可能含编号）：{caption}"
    try:
        from openai import OpenAI

        from src.services.http_client import create_openai_http_client

        client = OpenAI(
            api_key=api_key,
            base_url=base_url,
            timeout=p.visual_timeout,
            max_retries=1,
            http_client=create_openai_http_client(p.visual_timeout),
        )
        resp = client.chat.completions.create(
            model=p.visual_model,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {"type": "image_url", "image_url": {"url": data_url}},
                    ],
                }
            ],
            temperature=0.2,
            max_tokens=200,
        )
        text = (resp.choices[0].message.content or "").strip()
        if text:
            logger.info(
                "视觉摘要完成: %s (第%s页)", Path(image_path).name, page or "?"
            )
            return text
    except Exception as e:
        logger.warning("视觉摘要失败，图片按占位文本处理: %s", e)
    return None


def _parse_pdf_mineru(
    path: str,
    *,
    cmd: str = "mineru",
    timeout: int = 0,
) -> ParsedDocument | None:
    """子进程调用 MinerU CLI，读取输出 Markdown/JSON 重组为 ParsedDocument。

    timeout<=0 表示不限时；任何失败返回 None，由调用方回退现有链路。
    """
    if not _mineru_available(cmd):
        return None
    tmp = Path(tempfile.mkdtemp(prefix="mineru_"))
    try:
        logger.info(
            "MinerU 解析开始: %s (timeout=%s)", Path(path).name, timeout or "无"
        )
        proc = subprocess.run(
            [cmd, "-p", str(path), "-o", str(tmp)],
            capture_output=True,
            text=True,
            timeout=None if timeout <= 0 else timeout,
            encoding="utf-8",
            errors="replace",
        )
        if proc.returncode != 0:
            logger.warning(
                "MinerU 退出码 %s: %s",
                proc.returncode,
                (proc.stderr or proc.stdout).strip()[:300],
            )
            return None

        md_files = sorted(tmp.rglob("*.md"))
        json_files = sorted(tmp.rglob("*.json"))
        md_doc = None
        if md_files:
            md_doc = _mineru_md_to_pages(
                md_files[0].read_text(encoding="utf-8", errors="replace")
            )
        for jf in json_files:
            try:
                content = json.loads(jf.read_text(encoding="utf-8", errors="replace"))
            except Exception:
                continue
            if isinstance(content, dict) and isinstance(content.get("content_list"), list):
                doc = _mineru_json_to_document(content, base_dir=jf.parent)
                if doc.pages:
                    logger.info(
                        "MinerU 解析完成: %s, %d 页 / %d 块 (json)",
                        Path(path).name,
                        len(doc.pages),
                        len(doc.blocks),
                    )
                    return doc
        if md_doc and md_doc.pages:
            logger.info("MinerU 解析完成: %s, 回退 Markdown", Path(path).name)
            return md_doc
        logger.warning("MinerU 未产出可用文本: %s", Path(path).name)
        return None
    except subprocess.TimeoutExpired:
        logger.warning("MinerU 解析超时（%s 秒）: %s", timeout, Path(path).name)
        return None
    except Exception as e:
        logger.warning("MinerU 解析异常: %s", e)
        return None
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


def _parse_pdf(path: str) -> ParsedDocument:
    p = config.parsing
    pdf_kind = _pdf_kind(path)
    logger.info("PDF 类型判断: %s -> %s", Path(path).name, pdf_kind)

    if p.pdf_parser in ("mineru", "auto"):
        want_mineru = p.pdf_parser == "mineru" or pdf_kind in ("scanned", "mixed")
        if want_mineru:
            if _mineru_available(p.mineru_cmd):
                doc = _parse_pdf_mineru(
                    path, cmd=p.mineru_cmd, timeout=p.mineru_timeout
                )
                if doc and doc.full_text.strip():
                    return doc
            elif p.pdf_parser == "mineru":
                logger.warning(
                    "PDF_PARSER=mineru 但找不到命令 %s，回退现有链路", p.mineru_cmd
                )

    # 原生文本 PDF：直接提取
    try:
        doc = _parse_pdf_pymupdf4llm(path)
        if doc and doc.full_text.strip():
            return doc
        if p.pdf_use_ocr and not p.pdf_force_ocr:
            logger.info("PDF 空文本，OCR 重试: %s", Path(path).name)
            doc = _parse_pdf_pymupdf4llm(path, force_ocr=True)
            if doc and doc.full_text.strip():
                return doc
    except Exception as e:
        logger.warning("pymupdf4llm 失败，回退 fitz: %s", e)

    doc = _parse_pdf_fitz(path)
    if doc.full_text.strip():
        return doc

    if p.pdf_use_ocr and not p.pdf_force_ocr:
        try:
            ocr = _parse_pdf_pymupdf4llm(path, force_ocr=True)
            if ocr and ocr.full_text.strip():
                return ocr
        except Exception as e:
            logger.warning("PDF OCR 失败: %s", e)
    return doc


def _parse_txt(path: str) -> str:
    data = Path(path).read_bytes()
    for enc in ("utf-8-sig", "utf-8", "gbk"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    # ponytail: 怪编码替换非法字节，不阻断
    return data.decode("utf-8", errors="replace")


def _parse_plain(path: str) -> ParsedDocument:
    text = _parse_txt(path).strip()
    return ParsedDocument([ParsedPage(page=None, text=text)] if text else [])


def _table_to_text(table) -> str:
    rows: list[str] = []
    for row in table.rows:
        seen: set[int] = set()
        cells: list[str] = []
        for cell in row.cells:
            key = id(cell._tc)
            if key in seen:
                continue
            seen.add(key)
            if t := cell.text.strip():
                cells.append(t)
            for nested in cell.tables:
                if nt := _table_to_text(nested):
                    cells.append(nt)
        if cells:
            rows.append(" | ".join(cells))
    return "\n".join(rows)


def _iter_block_items(parent):
    from docx.document import Document as DocxDocument
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    body = parent.element.body if isinstance(parent, DocxDocument) else parent._element
    for child in body.iterchildren():
        if child.tag == qn("w:p"):
            yield Paragraph(child, parent)
        elif child.tag == qn("w:tbl"):
            yield Table(child, parent)


def _story_parts(container) -> list[str]:
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    parts: list[str] = []
    for block in _iter_block_items(container):
        if isinstance(block, Paragraph):
            if t := (block.text or "").strip():
                parts.append(t)
        elif isinstance(block, Table):
            if t := _table_to_text(block):
                parts.append(t)
    return parts


def _header_footer_parts(doc) -> list[str]:
    # 单节异常跳过，避免整份 docx 失败
    seen: set[str] = set()
    out: list[str] = []
    for section in doc.sections:
        for part in (section.header, section.footer):
            try:
                block = "\n".join(_story_parts(part)).strip()
            except Exception:
                continue
            if block and block not in seen:
                seen.add(block)
                out.append(block)
    return out


def _textbox_parts(doc) -> list[str]:
    from docx.oxml.ns import qn

    parts: list[str] = []
    for txbx in doc.element.body.iter(qn("w:txbxContent")):
        texts = [(n.text or "").strip() for n in txbx.iter(qn("w:t")) if (n.text or "").strip()]
        if texts:
            parts.append("\n".join(texts))
    return parts


def _parse_docx(path: str) -> ParsedDocument:
    from docx import Document

    doc = Document(path)
    parts = _header_footer_parts(doc) + _story_parts(doc) + _textbox_parts(doc)
    text = "\n\n".join(parts)
    return ParsedDocument([ParsedPage(page=None, text=text)] if text else [])


def _find_soffice() -> str | None:
    for name in ("soffice", "libreoffice"):
        if found := shutil.which(name):
            return found
    if os.name == "nt":
        for exe in ("soffice.com", "soffice.exe"):
            for pattern in (
                rf"C:\Program Files\LibreOffice\program\{exe}",
                rf"C:\Program Files (x86)\LibreOffice\program\{exe}",
            ):
                if Path(pattern).is_file():
                    return pattern
    return None


def _convert_doc_to_docx_soffice(path: str) -> Path | None:
    soffice = _find_soffice()
    if not soffice:
        return None

    out_dir = Path(tempfile.mkdtemp(prefix="exam_doc_"))
    dest: Path | None = None
    try:
        proc, _stdout, stderr = _run_capture(
            [soffice, "--headless", "--convert-to", "docx", "--outdir", str(out_dir), path],
            timeout=120,
        )
        if proc.returncode != 0:
            logger.warning(
                "LibreOffice 退出码 %s: %s", proc.returncode, stderr.strip()[:300]
            )
            return None
        converted = out_dir / f"{Path(path).stem}.docx"
        if not converted.is_file():
            return None
        fd, dest_name = tempfile.mkstemp(suffix=".docx", prefix="exam_doc_")
        os.close(fd)
        dest = Path(dest_name)
        shutil.copy2(converted, dest)
        return dest
    except (subprocess.CalledProcessError, subprocess.TimeoutExpired, OSError) as e:
        logger.warning("LibreOffice 转换失败: %s", e)
        if dest is not None:
            dest.unlink(missing_ok=True)
        return None
    finally:
        shutil.rmtree(out_dir, ignore_errors=True)


def _convert_doc_to_docx_word(path: str) -> Path | None:
    """Windows：无 LibreOffice 时用本机 Word COM 另存为 .docx。"""
    if os.name != "nt":
        return None
    src = str(Path(path).resolve())
    fd, dest_name = tempfile.mkstemp(suffix=".docx", prefix="exam_doc_")
    os.close(fd)
    dest = Path(dest_name)
    dest.unlink(missing_ok=True)
    src_ps = src.replace("'", "''")
    dest_ps = str(dest.resolve()).replace("'", "''")
    ps = (
        "$ErrorActionPreference='Stop'; "
        "$word=New-Object -ComObject Word.Application; "
        "$word.Visible=$false; "
        f"$doc=$word.Documents.Open('{src_ps}'); "
        f"$doc.SaveAs([ref]'{dest_ps}',[ref]16); "
        "$doc.Close(); $word.Quit(); "
        "[System.Runtime.InteropServices.Marshal]::ReleaseComObject($word)|Out-Null"
    )
    try:
        proc, _stdout, stderr = _run_capture(
            ["powershell", "-NoProfile", "-ExecutionPolicy", "Bypass", "-Command", ps],
            timeout=120,
        )
        if proc.returncode != 0:
            logger.warning(
                "Word COM 退出码 %s: %s", proc.returncode, stderr.strip()[:300]
            )
            return None
        if dest.is_file() and dest.stat().st_size > 0:
            return dest
        dest.unlink(missing_ok=True)
        return None
    except (subprocess.CalledProcessError, subprocess.TimeoutExpired, OSError) as e:
        logger.warning("Word COM 转换失败: %s", e)
        dest.unlink(missing_ok=True)
        return None


def _convert_doc_to_docx(path: str) -> Path | None:
    return _convert_doc_to_docx_soffice(path) or _convert_doc_to_docx_word(path)


def _parse_doc(path: str) -> ParsedDocument:
    docx_tmp = _convert_doc_to_docx(path)
    if docx_tmp:
        try:
            return _parse_docx(str(docx_tmp))
        finally:
            docx_tmp.unlink(missing_ok=True)
    raise BadRequestException(
        "无法解析 .doc 文件。请安装 LibreOffice 或 Microsoft Word，或将文件另存为 .docx 后重试。"
    )


def _shape_texts(shape) -> list[str]:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        return [t for child in shape.shapes for t in _shape_texts(child)]

    texts: list[str] = []
    if getattr(shape, "has_table", False):
        rows = []
        for row in shape.table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            texts.append("\n".join(rows))
    if getattr(shape, "has_text_frame", False):
        if t := shape.text_frame.text.strip():
            texts.append(t)
    elif hasattr(shape, "text") and (t := (shape.text or "").strip()):
        texts.append(t)
    return texts


def _convert_pptx_to_pdf_soffice(path: str) -> Path | None:
    """LibreOffice headless 将 PPTX 转为 PDF，可覆盖 SmartArt/图表/母版文本。"""
    soffice = _find_soffice()
    if not soffice:
        return None

    out_dir = Path(tempfile.mkdtemp(prefix="exam_pptx_"))
    profile_dir = Path(tempfile.mkdtemp(prefix="exam_lo_profile_"))
    dest: Path | None = None
    try:
        user_install = "file:///" + str(profile_dir).replace("\\", "/")
        proc, _stdout, stderr = _run_capture(
            [
                soffice,
                "--headless",
                "--norestore",
                "--nolockcheck",
                "--convert-to",
                "pdf",
                "--outdir",
                str(out_dir),
                f"-env:UserInstallation={user_install}",
                path,
            ],
            timeout=180,
        )
        if proc.returncode != 0:
            logger.warning(
                "LibreOffice 退出码 %s: %s", proc.returncode, stderr.strip()[:300]
            )
            return None
        converted = out_dir / f"{Path(path).stem}.pdf"
        if not converted.is_file():
            return None
        fd, dest_name = tempfile.mkstemp(suffix=".pdf", prefix="exam_pptx_")
        os.close(fd)
        dest = Path(dest_name)
        shutil.copy2(converted, dest)
        return dest
    except (subprocess.CalledProcessError, subprocess.TimeoutExpired, OSError) as e:
        logger.warning("LibreOffice PPTX 转 PDF 失败: %s", e)
        if dest is not None:
            dest.unlink(missing_ok=True)
        return None
    finally:
        shutil.rmtree(out_dir, ignore_errors=True)
        shutil.rmtree(profile_dir, ignore_errors=True)


def _convert_pptx_to_pdf_powerpoint(path: str) -> Path | None:
    """Windows：无 LibreOffice 时用本机 PowerPoint COM 导出 PDF（ppSaveAsPDF=32）。"""
    if os.name != "nt":
        return None
    src = str(Path(path).resolve())
    fd, dest_name = tempfile.mkstemp(suffix=".pdf", prefix="exam_pptx_")
    os.close(fd)
    dest = Path(dest_name)
    dest.unlink(missing_ok=True)
    src_ps = src.replace("'", "''")
    dest_ps = str(dest.resolve()).replace("'", "''")
    ps = (
        "$ErrorActionPreference='Stop'; "
        "$pp=New-Object -ComObject PowerPoint.Application; "
        "try { "
        f"$pres=$pp.Presentations.Open('{src_ps}',-1,0,0); "
        f"$pres.SaveAs([ref]'{dest_ps}',[ref]32); "
        "$pres.Close(); "
        "} finally { "
        "$pp.Quit(); "
        "[System.Runtime.InteropServices.Marshal]::ReleaseComObject($pp)|Out-Null "
        "}"
    )
    try:
        subprocess.run(
            ["powershell", "-NoProfile", "-ExecutionPolicy", "Bypass", "-Command", ps],
            check=True,
            capture_output=True,
            text=True,
            timeout=180,
        )
        if dest.is_file() and dest.stat().st_size > 0:
            return dest
        dest.unlink(missing_ok=True)
        return None
    except (subprocess.CalledProcessError, subprocess.TimeoutExpired, OSError) as e:
        logger.warning("PowerPoint COM 导出 PDF 失败: %s", e)
        dest.unlink(missing_ok=True)
        return None


def _convert_pptx_to_pdf(path: str) -> Path | None:
    return _convert_pptx_to_pdf_soffice(path) or _convert_pptx_to_pdf_powerpoint(path)


def _parse_pptx_direct(path: str) -> ParsedDocument:
    from pptx import Presentation

    pages: list[ParsedPage] = []
    for i, slide in enumerate(Presentation(path).slides, 1):
        texts = [t for shape in slide.shapes for t in _shape_texts(shape)]
        if slide.has_notes_slide:
            if notes := slide.notes_slide.notes_text_frame.text.strip():
                texts.append(f"[备注]\n{notes}")
        if texts:
            pages.append(ParsedPage(page=i, text="\n".join(texts)))
    return ParsedDocument(pages)


def _tesseract_available() -> bool:
    if shutil.which("tesseract"):
        return True
    if os.name == "nt" and Path(r"C:\Program Files\Tesseract-OCR\tesseract.exe").is_file():
        return True
    return False


def _clean_pdf_markdown(text: str) -> str:
    """清理 pymupdf4llm 输出的 markdown 杂质（图片占位符、标题标记等）。"""
    text = re.sub(r"\*\*==>.*?<==\*\*", "", text, flags=re.S)
    text = re.sub(r"^\s*#+\s*$", "", text, flags=re.M)
    text = re.sub(r"^\s*#{1,6}\s+", "", text, flags=re.M)
    text = text.replace("**", "").replace("__", "")
    text = re.sub(r"^\s*\|?[\s:|-]+\|?\s*$", "", text, flags=re.M)
    lines = [ln.rstrip() for ln in text.splitlines() if ln.strip()]
    return "\n".join(lines).strip()


def _merge_pages_by_length(a: list[ParsedPage], b: list[ParsedPage]) -> ParsedDocument:
    """同一 PPT 的两种提取结果按页择优合并（PDF 页号与 PPT 页号一致）。"""
    pages: list[ParsedPage] = []
    for i in range(max(len(a), len(b))):
        pa = a[i] if i < len(a) else None
        pb = b[i] if i < len(b) else None
        if pa and pb:
            pages.append(pa if len(pa.text) >= len(pb.text) else pb)
        else:
            pages.append(pa or pb)
    return ParsedDocument([p for p in pages if p.text and p.text.strip()])


def _parse_pptx(path: str) -> ParsedDocument:
    """优先转 PDF 走 pymupdf4llm 管线，清理杂质后与 python-pptx 按页择优合并；
    图片型 PPT 自动尝试 OCR（需系统安装 Tesseract）。"""
    direct = _parse_pptx_direct(path)
    pdf_tmp = _convert_pptx_to_pdf(path)
    if not pdf_tmp:
        return direct
    try:
        doc = _parse_pdf(str(pdf_tmp))
        pdf_pages = [
            ParsedPage(page=p.page, text=_clean_pdf_markdown(p.text))
            for p in doc.pages
            if _clean_pdf_markdown(p.text).strip()
        ]
        direct_chars = sum(len(p.text) for p in direct.pages)
        pdf_chars = sum(len(p.text) for p in pdf_pages)
        if pdf_chars < max(direct_chars, 1) * 0.5 and _tesseract_available():
            try:
                ocr_doc = _parse_pdf_pymupdf4llm(str(pdf_tmp), force_ocr=True)
                ocr_pages = [
                    ParsedPage(page=p.page, text=_clean_pdf_markdown(p.text))
                    for p in ocr_doc.pages
                    if _clean_pdf_markdown(p.text).strip()
                ]
                ocr_chars = sum(len(p.text) for p in ocr_pages)
                if ocr_chars > pdf_chars:
                    logger.info("PPTX 图片 OCR 生效: %s -> %d 字", Path(path).name, ocr_chars)
                    pdf_pages = ocr_pages
            except Exception as e:
                logger.warning("PPTX 图片 OCR 失败: %s", e)
        return _merge_pages_by_length(direct.pages, pdf_pages)
    except Exception as e:
        logger.warning("PPTX 转 PDF 解析失败，回退 python-pptx: %s", e)
        return direct
    finally:
        pdf_tmp.unlink(missing_ok=True)


_PARSERS = {
    ".pdf": _parse_pdf,
    ".txt": _parse_plain,
    ".md": _parse_plain,
    ".doc": _parse_doc,
    ".docx": _parse_docx,
    ".pptx": _parse_pptx,
}


def parse_file(path: str) -> ParsedDocument:
    ext = Path(path).suffix.lower()
    if ext not in _PARSERS:
        raise UnsupportedFormatException(
            f"不支持的文件格式: {ext}，仅接受 PDF/TXT/MD/DOC/DOCX/PPTX"
        )
    if not os.path.exists(path):
        raise BadRequestException(f"文件不存在: {path}")

    try:
        doc = _PARSERS[ext](path)
    except (UnsupportedFormatException, BadRequestException):
        raise
    except Exception as e:
        logger.error("文件解析失败 (%s): %s", path, e)
        raise BadRequestException(f"文件解析失败: {e}") from e

    if not doc.full_text.strip():
        raise BadRequestException("文件内容为空，无法入库")

    logger.info(
        "解析完成: %s pages=%d chars=%d",
        Path(path).name,
        len(doc.pages),
        len(doc.full_text),
    )
    return doc
