"""上传解析：产品支持的格式（MD/TXT/DOCX/PPTX）+ 拒答空文件。"""

import pytest

from src.exceptions import BadRequestException, UnsupportedFormatException
from src.services.parsing import parse_file


class TestParsePlain:
    def test_parse_markdown(self, sample_md_file):
        doc = parse_file(sample_md_file)
        assert "测试文档" in doc.full_text
        assert len(doc.pages) == 1

    def test_parse_txt(self, sample_txt_file):
        assert "纯文本" in parse_file(sample_txt_file).full_text

    def test_empty_and_unsupported(self, temp_dir):
        empty = temp_dir / "empty.txt"
        empty.write_text("   \n", encoding="utf-8")
        with pytest.raises(BadRequestException):
            parse_file(str(empty))
        bad = temp_dir / "data.bin"
        bad.write_bytes(b"\x00\x01")
        with pytest.raises(UnsupportedFormatException):
            parse_file(str(bad))


class TestParseOffice:
    def test_parse_docx(self, temp_dir):
        from docx import Document

        path = temp_dir / "notes.docx"
        doc = Document()
        doc.add_heading("傅里叶变换", level=1)
        doc.add_paragraph("连续时间傅里叶变换的定义。")
        doc.save(str(path))
        parsed = parse_file(str(path))
        assert "傅里叶变换" in parsed.full_text
        assert "连续时间" in parsed.full_text

    def test_parse_pptx(self, temp_dir):
        from pptx import Presentation

        path = temp_dir / "slides.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "采样定理"
        slide.placeholders[1].text = "奈奎斯特频率 fs/2"
        prs.save(str(path))
        parsed = parse_file(str(path))
        assert "采样定理" in parsed.full_text
        assert parsed.pages[0].page == 1

    def test_parse_pptx_fallback_without_converter(self, temp_dir, monkeypatch):
        """无 LibreOffice/PowerPoint 时必须回退 python-pptx 提取。"""
        from pptx import Presentation
        from src.services import parsing

        path = temp_dir / "fallback.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "回退测试"
        prs.save(str(path))

        monkeypatch.setattr(parsing, "_convert_pptx_to_pdf", lambda p: None)
        parsed = parse_file(str(path))
        assert "回退测试" in parsed.full_text
        assert parsed.pages[0].page == 1


class TestMineruMapping:
    """MinerU 输出重组逻辑（子进程调用本身不在此测）。"""

    def test_json_to_pages_groups_by_page(self):
        from src.services import parsing

        content = {
            "content_list": [
                {"type": "text", "page_idx": 0, "text": "第一章 引言"},
                {"type": "formula", "page_idx": 0, "latex": "E=mc^2"},
                {"type": "title", "page_idx": 1, "text": "1.1 通信系统模型"},
                {
                    "type": "table",
                    "page_idx": 1,
                    "html": "<table><tr><td>信源</td><td>编码</td></tr></table>",
                },
                {"type": "image", "page_idx": 1},
            ]
        }
        doc = parsing._mineru_json_to_pages(content)
        assert len(doc.pages) == 2
        assert doc.pages[0].page == 1
        assert "第一章" in doc.pages[0].text
        assert "$$E=mc^2$$" in doc.pages[0].text
        assert "信源 | 编码" in doc.pages[1].text

    def test_md_fallback_single_page(self):
        from src.services import parsing

        doc = parsing._mineru_md_to_pages("## 标题\n正文")
        assert len(doc.pages) == 1
        assert doc.pages[0].page is None
        assert "正文" in doc.full_text

    def test_html_table_to_text(self):
        from src.services import parsing

        out = parsing._html_table_to_text(
            "<table><tr><td>a</td><td>b</td></tr><tr><td>1</td><td>2</td></tr></table>"
        )
        assert "a | b" in out
        assert "1 | 2" in out


class TestMineruStructure:
    """结构还原：block 类型、标题层级、表格 HTML、图片路径与说明关联。"""

    def _content(self, tmp_path):
        img = tmp_path / "fig1.png"
        img.write_bytes(b"\x89PNG\r\n\x1a\n" + b"\x00" * 16)
        return {
            "content_list": [
                {"type": "header", "page_idx": 0, "text": "信号与系统"},
                {"type": "title", "page_idx": 0, "text": "第一章 绪论", "level": 1},
                {"type": "title", "page_idx": 0, "text": "1.1 通信系统模型", "level": 2},
                {"type": "text", "page_idx": 0, "text": "通信系统由信源、信道和信宿组成。"},
                {
                    "type": "table",
                    "page_idx": 1,
                    "html": (
                        "<table><tr><th>名称</th><th>说明</th></tr>"
                        "<tr><td>信源</td><td>信息发起点</td></tr></table>"
                    ),
                },
                {"type": "image", "page_idx": 1, "img_path": "fig1.png"},
                {"type": "image_caption", "page_idx": 1, "text": "图1-1 通信系统框图"},
                {"type": "formula", "page_idx": 1, "latex": "C = B\\log_2(1+SNR)"},
            ]
        }

    def test_structure_restored(self, tmp_path):
        from src.services import parsing

        doc = parsing._mineru_json_to_document(
            self._content(tmp_path), base_dir=tmp_path
        )
        assert doc.has_blocks
        # header + title×2 + text + table + image + formula = 7 块（caption 挂载不单独成块）
        assert len(doc.blocks) == 7

        titles = [b for b in doc.blocks if b.block_type == "title"]
        assert [t.level for t in titles] == [1, 2]
        assert titles[0].page == 1 and titles[1].page == 1

        table = next(b for b in doc.blocks if b.block_type == "table")
        assert table.html and "名称" in table.html
        assert parsing._table_headers(table.html) == "名称 | 说明"
        assert table.page == 2

        image = next(b for b in doc.blocks if b.block_type == "image")
        assert image.image_path == str(tmp_path / "fig1.png")
        assert image.caption == "图1-1 通信系统框图"

        formula = next(b for b in doc.blocks if b.block_type == "formula")
        assert "$$C = B\\log_2(1+SNR)$$" in formula.text

        header = next(b for b in doc.blocks if b.block_type == "header")
        assert header.text == "信号与系统"
        assert len(doc.pages) == 2

    def test_table_headers_extraction(self):
        from src.services import parsing

        assert (
            parsing._table_headers(
                "<table><tr><th>信源</th><th>编码</th></tr><tr><td>a</td><td>b</td></tr></table>"
            )
            == "信源 | 编码"
        )
        assert parsing._table_headers("<table><tr><td>a</td></tr></table>") == ""
