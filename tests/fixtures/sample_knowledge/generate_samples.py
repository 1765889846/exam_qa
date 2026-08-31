"""一次性生成多格式测试讲义 → data/knowledge/（可重复运行覆盖）。"""

from __future__ import annotations

from pathlib import Path
import sys

OUT = Path(__file__).resolve().parents[3] / "data" / "knowledge"


def write_txt() -> None:
    (OUT / "laplace-transform.txt").write_text(
        """第5章 拉普拉斯变换

一、拉普拉斯变换的定义

定义：设函数 f(t) 在 t>=0 上有定义，若积分
  F(s) = ∫_0^∞ f(t) e^{-st} dt
收敛，则称 F(s) 为 f(t) 的拉普拉斯变换，记作 L{f(t)}=F(s)。

二、常用公式

1. L{1} = 1/s  (Re(s)>0)
2. L{t^n} = n!/s^{n+1}
3. L{e^{at}} = 1/(s-a)
4. 微分性质：L{f'(t)} = s F(s) - f(0)

三、例题

例：求 f(t)=e^{-2t} 的拉普拉斯变换。
解：L{e^{-2t}} = 1/(s+2)，其中 Re(s)>-2。

例：用微分性质求 L{cos(ωt)}。
已知 L{sin(ωt)}=ω/(s^2+ω^2)，对 t 微分并整理可得
L{cos(ωt)}=s/(s^2+ω^2)。
""",
        encoding="utf-8",
    )


def write_docx() -> None:
    from docx import Document

    doc = Document()
    doc.add_heading("第2章 特征值与特征向量", level=1)
    doc.add_heading("特征值的定义", level=2)
    doc.add_paragraph(
        "定义：设 A 为 n 阶方阵，若存在数 λ 与非零向量 x，使得 A x = λ x，"
        "则称 λ 为 A 的特征值，x 为属于 λ 的特征向量。"
    )
    doc.add_heading("特征方程（公式）", level=2)
    doc.add_paragraph("特征多项式 det(A − λI)=0 的根即为特征值。")
    doc.add_paragraph(
        "对 2×2 矩阵 A=[[a,b],[c,d]]，特征方程为 λ² − (a+d)λ + (ad−bc)=0。"
    )
    doc.add_heading("例题", level=2)
    doc.add_paragraph("例：求 A=[[2,1],[1,2]] 的特征值与特征向量。")
    doc.add_paragraph(
        "解：det(A−λI)=(2−λ)²−1=λ²−4λ+3=0，得 λ₁=3，λ₂=1。"
        "对 λ=3，(A−3I)x=0 得 x=k(1,1)^T；对 λ=1，得 x=k(1,−1)^T。"
    )
    doc.save(OUT / "eigenvalues.docx")


def write_pptx() -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_slide(title: str, lines: list[str]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        body.clear()
        for i, line in enumerate(lines):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(20)

    add_slide(
        "第6章 贝叶斯公式",
        [
            "本章知识点：条件概率、贝叶斯定理、全概率公式",
            "用于溯知 PPTX 入库与章节概览测试",
        ],
    )
    add_slide(
        "定义：条件概率",
        [
            "P(A|B)=P(A∩B)/P(B)，其中 P(B)>0",
            "含义：在事件 B 已发生的条件下，A 发生的概率",
        ],
    )
    add_slide(
        "公式：贝叶斯定理",
        [
            "P(A_i|B)=P(B|A_i)P(A_i)/Σ_j P(B|A_j)P(A_j)",
            "先验 P(A_i) → 似然 P(B|A_i) → 后验 P(A_i|B)",
        ],
    )
    add_slide(
        "例题",
        [
            "某病患病率 1%，检测真阳性率 99%，假阳性率 5%。",
            "若某人检测呈阳性，患病后验概率约为多少？",
            "解：P(病|阳)=0.99×0.01/(0.99×0.01+0.05×0.99)≈0.167（约 16.7%）",
        ],
    )
    prs.save(OUT / "bayes-theorem.pptx")


def write_pdf() -> None:
    import fitz

    text = """第7章 一阶微分方程

一、定义
形如 y' = f(x,y) 的方程称为一阶常微分方程。
若可写成 M(x,y)dx + N(x,y)dy = 0，则称全微分形式。

二、公式（可分离变量）
若 y' = g(x)h(y)，且 h(y)≠0，则
  ∫ dy/h(y) = ∫ g(x) dx + C

三、例题
例：求解初值问题 y' = 2xy，y(0)=1。
解：分离变量 dy/y = 2x dx（y≠0），
积分得 ln|y| = x² + C₁，故 y = C e^{x²}。
由 y(0)=1 得 C=1，因此 y = e^{x²}。

四、要点
初值问题解的存在唯一性依赖 f 对 y 的利普希茨条件。
"""
    pdf = fitz.open()
    page = pdf.new_page(width=595, height=842)
    # china-s：PyMuPDF 内置宋体子集，中文可检索
    page.insert_textbox(
        fitz.Rect(50, 50, 545, 792),
        text,
        fontname="china-s",
        fontsize=12,
        align=0,
    )
    pdf.save(OUT / "ode-first-order.pdf")
    pdf.close()


def write_doc_via_word() -> None:
    """生成真正的 .doc（Word COM）。失败则写同名说明 txt 提示装 LibreOffice。"""
    import tempfile

    content = """第8章 行列式

一、定义
n 阶行列式 det(A) 是按某一行（列）展开得到的标量。
二阶行列式 |a b; c d| = ad − bc。

二、公式
1. 行列互换，行列式变号。
2. 某行（列）公因子可提出。
3. 三角形行列式等于主对角元之积。
4. 克拉默法则：Ax=b（det A≠0）时，x_i = det(A_i)/det(A)。

三、例题
例：计算 D = |1 2 3; 0 4 5; 0 0 6|。
解：上三角，D = 1×4×6 = 24。

例：用克拉默法则解 {2x+y=5, x−y=1}。
系数行列式 Δ=|2 1; 1 −1|=−3，Δ_x=|5 1; 1 −1|=−6，Δ_y=|2 5; 1 1|=−3，
故 x=2，y=1。
"""
    # Word COM（无 pywin32）
    ps = r"""
$ErrorActionPreference = 'Stop'
$path = $args[0]
$text = Get-Content -Raw -Encoding UTF8 $args[1]
$word = New-Object -ComObject Word.Application
$word.Visible = $false
$doc = $word.Documents.Add()
$doc.Content.Text = $text
$doc.SaveAs([ref]$path, [ref]0)
$doc.Close()
$word.Quit()
[System.Runtime.InteropServices.Marshal]::ReleaseComObject($word) | Out-Null
"""
    with tempfile.NamedTemporaryFile(
        "w", suffix=".ps1", delete=False, encoding="utf-8"
    ) as sf:
        sf.write(ps)
        ps_path = sf.name
    with tempfile.NamedTemporaryFile(
        "w", suffix=".txt", delete=False, encoding="utf-8"
    ) as tf:
        tf.write(content)
        txt_path = tf.name

    doc_path = OUT / "determinants.doc"
    import subprocess

    try:
        subprocess.run(
            [
                "powershell",
                "-NoProfile",
                "-ExecutionPolicy",
                "Bypass",
                "-File",
                ps_path,
                str(doc_path.resolve()),
                txt_path,
            ],
            check=True,
            capture_output=True,
            text=True,
            timeout=120,
        )
        if not doc_path.is_file():
            raise RuntimeError("Word SaveAs 未生成 .doc")
    finally:
        Path(ps_path).unlink(missing_ok=True)
        Path(txt_path).unlink(missing_ok=True)


SAMPLE_FILES = (
    "laplace-transform.txt",
    "eigenvalues.docx",
    "bayes-theorem.pptx",
    "ode-first-order.pdf",
    "determinants.doc",
)


def cleanup() -> None:
    """删除 generate_samples 写入的测试样本。"""
    for name in SAMPLE_FILES:
        (OUT / name).unlink(missing_ok=True)


def main() -> None:
    if "--clean" in sys.argv:
        cleanup()
        print("已清理测试样本")
        return
    OUT.mkdir(parents=True, exist_ok=True)
    write_txt()
    print("txt")
    write_docx()
    print("docx")
    write_pptx()
    print("pptx")
    write_pdf()
    print("pdf")
    write_doc_via_word()
    print("doc")
    for p in sorted(OUT.iterdir()):
        if p.is_file():
            print(f"  {p.name:40s} {p.stat().st_size:6d} B")


if __name__ == "__main__":
    main()
